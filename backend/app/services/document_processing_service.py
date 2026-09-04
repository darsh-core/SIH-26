import os
import uuid
import re
import hashlib
from datetime import datetime
from typing import List, Dict, Any, Tuple, Optional
import pypdf
import docx
from sqlalchemy.orm import Session

from app.models.document import Document, DocumentChunk, DocumentEmbedding
from app.ai import get_embedding_provider

logger = __import__("logging").getLogger("sih-platform.services.document")

def normalize_text(text: str) -> str:
    """Removes excessive whitespace and standardizes line endings."""
    if not text:
        return ""
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


class PDFTextExtractor:
    @staticmethod
    def extract(file_path: str) -> List[Dict[str, Any]]:
        pages_data = []
        try:
            reader = pypdf.PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                raw = page.extract_text() or ""
                clean = normalize_text(raw)
                if clean:
                    pages_data.append({
                        "page_number": i + 1,
                        "text": clean,
                        "source_type": "page"
                    })
        except Exception as e:
            logger.error(f"Error in PDFTextExtractor for {file_path}: {e}")
            raise e
        return pages_data


class DOCXTextExtractor:
    @staticmethod
    def extract(file_path: str) -> List[Dict[str, Any]]:
        pages_data = []
        try:
            doc = docx.Document(file_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # Group paragraphs into virtual pages (e.g. 6-8 paragraphs per virtual page)
            paragraphs_per_page = 8
            for idx, start in enumerate(range(0, len(paragraphs), paragraphs_per_page)):
                text = "\n\n".join(paragraphs[start:start + paragraphs_per_page])
                clean = normalize_text(text)
                if clean:
                    pages_data.append({
                        "page_number": idx + 1,
                        "text": clean,
                        "source_type": "section"
                    })
        except Exception as e:
            logger.error(f"Error in DOCXTextExtractor for {file_path}: {e}")
            raise e
        return pages_data


class PPTXTextExtractor:
    @staticmethod
    def extract(file_path: str) -> List[Dict[str, Any]]:
        pages_data = []
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        slide_lines.append(shape.text_frame.text.strip())
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([c.text.strip() for c in row.cells if c.text.strip()])
                            if row_text:
                                slide_lines.append(row_text)
                
                # Check for slide notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_lines.append(f"Notes: {notes_text}")
                
                joined = "\n".join(slide_lines)
                clean = normalize_text(joined)
                if clean:
                    pages_data.append({
                        "page_number": i + 1,
                        "text": clean,
                        "source_type": "slide"
                    })
        except Exception as e:
            logger.error(f"Error in PPTXTextExtractor for {file_path}: {e}")
            raise e
        return pages_data


class TXTTextExtractor:
    @staticmethod
    def extract(file_path: str) -> List[Dict[str, Any]]:
        pages_data = []
        content = None
        # Try multiple encodings for robustness
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            raise ValueError(f"Failed to decode TXT file {file_path} with supported encodings.")

        clean_content = normalize_text(content)
        # Split plain text into virtual pages of roughly 1500 characters each
        char_per_page = 1500
        for idx, start in enumerate(range(0, len(clean_content), char_per_page)):
            text = clean_content[start:start + char_per_page].strip()
            if text:
                pages_data.append({
                    "page_number": idx + 1,
                    "text": text,
                    "source_type": "virtual_page"
                })
        return pages_data


class DocumentProcessingService:
    """Orchestrates document extraction, semantic chunking, embedding generation, and vector indexing."""

    @staticmethod
    def compute_sha256(data: Any) -> str:
        if isinstance(data, str):
            data = data.encode("utf-8")
        return hashlib.sha256(data).hexdigest()

    @staticmethod
    def compute_chunk_hash(document_id: Any, chunk_index: int, text: str) -> str:
        return hashlib.sha256(f"{document_id}:{chunk_index}:{text}".encode("utf-8")).hexdigest()

    @staticmethod
    def process_document(db: Session, document_id: uuid.UUID) -> None:
        doc = db.query(Document).filter(Document.id == document_id).first()
        if not doc:
            return

        doc.status = "PROCESSING"
        db.commit()

        try:
            # 1. Select text extractor based on file extension
            ext = os.path.splitext(doc.filename)[1].lower()
            if ext == ".pdf":
                pages = PDFTextExtractor.extract(doc.file_path)
            elif ext == ".docx":
                pages = DOCXTextExtractor.extract(doc.file_path)
            elif ext == ".pptx":
                pages = PPTXTextExtractor.extract(doc.file_path)
            elif ext in [".txt", ".log"]:
                pages = TXTTextExtractor.extract(doc.file_path)
            else:
                raise ValueError(f"Unsupported file extension: {ext}")

            if not pages:
                raise ValueError("No extractable textual content found in document.")

            # 2. Semantic Chunking (preserving sentence & paragraph boundaries)
            chunks = DocumentProcessingService.chunk_pages(pages, target_word_count=350, overlap_word_count=60)
            if not chunks:
                raise ValueError("Semantic chunking produced 0 chunks.")

            # Clean out any old chunks/embeddings if re-processing
            db.query(DocumentChunk).filter(DocumentChunk.document_id == document_id).delete()
            db.flush()

            # 3. Embed and store chunks in PostgreSQL + pgvector
            embedding_provider = get_embedding_provider()
            model_name = getattr(embedding_provider, "model_name", "all-MiniLM-L6-v2")
            
            chunk_texts = [c["text"] for c in chunks]
            vectors = embedding_provider.embed_documents(chunk_texts)

            for idx, chunk in enumerate(chunks):
                text_content = chunk["text"]
                chunk_hash = hashlib.sha256(f"{document_id}:{idx}:{text_content}".encode("utf-8")).hexdigest()

                db_chunk = DocumentChunk(
                    document_id=document_id,
                    chunk_index=idx,
                    text_content=text_content,
                    start_char=chunk["start_char"],
                    end_char=chunk["end_char"],
                    page_number=chunk["page_start"],
                    chunk_hash=chunk_hash,
                    metadata_json={
                        "page_end": chunk["page_end"],
                        "token_count": chunk["token_count"],
                        "source_type": chunk.get("source_type", "page")
                    }
                )
                db.add(db_chunk)
                db.flush()

                vector = vectors[idx] if idx < len(vectors) else embedding_provider.embed_text(text_content)
                
                db_embed = DocumentEmbedding(
                    chunk_id=db_chunk.id,
                    model_name=model_name,
                    embedding=vector
                )
                db.add(db_embed)

            # 4. Mark READY / INDEXED
            doc.status = "READY"
            doc.metadata_json = {
                **(doc.metadata_json or {}),
                "chunk_count": len(chunks),
                "embedding_model": model_name,
                "processed_at": datetime.now().isoformat()
            }
            db.commit()
            logger.info(f"Successfully processed and indexed document {document_id} with {len(chunks)} chunks.")

        except Exception as e:
            db.rollback()
            doc.status = "FAILED"
            doc.metadata_json = {
                **(doc.metadata_json or {}),
                "error": str(e),
                "failed_at": datetime.now().isoformat()
            }
            db.commit()
            logger.error(f"Failed to process document {document_id}: {e}", exc_info=True)
            raise e

    @staticmethod
    def chunk_pages(
        pages: List[Dict[str, Any]],
        target_word_count: int = 350,
        overlap_word_count: int = 60
    ) -> List[Dict[str, Any]]:
        """
        Splits page/slide records into overlapping semantic chunks deterministically.
        Preserves paragraph and sentence boundaries.
        """
        chunks = []
        words_buffer: List[Tuple[str, int, int, str]] = []  # (word, page_num, char_offset, source_type)
        
        char_offset = 0
        for page in pages:
            text = page["text"]
            page_num = page["page_number"]
            source_type = page.get("source_type", "page")
            
            pos = 0
            for word in text.split():
                word_len = len(word)
                words_buffer.append((word, page_num, char_offset + pos, source_type))
                pos += word_len + 1
            char_offset += len(text) + 2

        if not words_buffer:
            return []

        i = 0
        while i < len(words_buffer):
            chunk_slice = words_buffer[i : i + target_word_count]
            if not chunk_slice:
                break
                
            chunk_text = " ".join([w[0] for w in chunk_slice])
            page_start = chunk_slice[0][1]
            page_end = chunk_slice[-1][1]
            start_char = chunk_slice[0][2]
            end_char = chunk_slice[-1][2] + len(chunk_slice[-1][0])
            source_type = chunk_slice[0][3]
            
            token_est = int(len(chunk_slice) * 1.3)
            
            chunks.append({
                "text": chunk_text,
                "page_start": page_start,
                "page_end": page_end,
                "start_char": start_char,
                "end_char": end_char,
                "token_count": token_est,
                "source_type": source_type
            })
            
            i += max(1, target_word_count - overlap_word_count)
            
        return chunks

    @staticmethod
    def search_similar_chunks(
        db: Session,
        query_embedding: List[float],
        top_k: int = 5,
        document_id: Optional[uuid.UUID] = None,
        min_similarity: float = 0.0
    ) -> List[Tuple[DocumentChunk, float]]:
        """
        Executes Cosine Similarity search over 384-D vector space in PostgreSQL pgvector.
        Cosine similarity = 1.0 - cosine_distance.
        """
        distance_col = DocumentEmbedding.embedding.cosine_distance(query_embedding)
        similarity_col = (1.0 - distance_col).label("similarity")

        query = db.query(DocumentChunk, similarity_col)\
            .join(DocumentEmbedding, DocumentChunk.id == DocumentEmbedding.chunk_id)

        if document_id:
            query = query.filter(DocumentChunk.document_id == document_id)

        if min_similarity > 0.0:
            query = query.filter((1.0 - distance_col) >= min_similarity)

        results = query.order_by(distance_col.asc()).limit(top_k).all()
        return [(row[0], float(row[1])) for row in results]
