import os
import tempfile
import uuid
import pytest
from pydantic import BaseModel
from sqlalchemy.orm import Session
from fastapi.testclient import TestClient

from app.core.config import settings
from app.ai.ollama_client import OllamaClient
from app.ai.llm import MockLLMProvider, OllamaLLMProvider, GroqLLMProvider
from app.ai.embeddings import MockEmbeddingProvider, SentenceTransformerEmbeddingProvider
from app.ai import get_llm_provider, get_embedding_provider
from app.services.document_processing_service import (
    PDFTextExtractor,
    DOCXTextExtractor,
    PPTXTextExtractor,
    TXTTextExtractor,
    DocumentProcessingService
)
from app.ai.validators import MCQValidator
from app.schemas.assessment import GeneratedMCQ, GeneratedMCQOption
from app.models.competency import Competency, CompetencyFramework, JobRole, RoleCompetency
from app.models.document import Document
from app.ai.role_assessment_generator import RoleDiagnosticGenerator


class SampleSchema(BaseModel):
    title: str
    count: int


def test_ai_provider_selection():
    """Verify that get_llm_provider returns configured provider."""
    original_provider = settings.AI_PROVIDER
    try:
        settings.AI_PROVIDER = "mock"
        os.environ["AI_PROVIDER"] = "mock"
        p_mock = get_llm_provider(force_refresh=True)
        assert isinstance(p_mock, MockLLMProvider)
        assert p_mock.provider_name == "mock"

        settings.AI_PROVIDER = "ollama"
        os.environ["AI_PROVIDER"] = "ollama"
        p_ollama = get_llm_provider(force_refresh=True)
        assert isinstance(p_ollama, OllamaLLMProvider)
        assert p_ollama.provider_name == "ollama"

        settings.AI_PROVIDER = "groq"
        os.environ["AI_PROVIDER"] = "groq"
        p_groq = get_llm_provider(force_refresh=True)
        assert isinstance(p_groq, GroqLLMProvider)
        assert p_groq.provider_name == "groq"
    finally:
        settings.AI_PROVIDER = original_provider
        os.environ["AI_PROVIDER"] = original_provider
        get_llm_provider(force_refresh=True)


def test_mock_llm_provider_structured():
    """Verify MockLLMProvider returns valid structured schema instances."""
    mock = MockLLMProvider()
    res = mock.generate_structured("Generate sample data", SampleSchema)
    assert isinstance(res, SampleSchema)
    assert isinstance(res.title, str)
    assert isinstance(res.count, int)


def test_ollama_client_properties():
    """Verify OllamaClient initialization and configuration."""
    client = OllamaClient(
        base_url="http://localhost:11434",
        model="llama3.2:latest",
        timeout=10
    )
    assert client.base_url == "http://localhost:11434"
    assert client.model == "llama3.2:latest"
    assert client.timeout == 10
    is_avail = client.is_available()
    assert isinstance(is_avail, bool)


def test_embedding_dimensions_canonical_384():
    """Verify both Mock and SentenceTransformer providers output canonical 384 dimensions."""
    mock_emb = MockEmbeddingProvider()
    v1 = mock_emb.embed_text("Sample survey frame")
    assert len(v1) == 384
    assert mock_emb.dimension == 384

    docs = ["Doc 1 text", "Doc 2 text"]
    v_docs = mock_emb.embed_documents(docs)
    assert len(v_docs) == 2
    assert len(v_docs[0]) == 384
    assert len(v_docs[1]) == 384

    try:
        st_emb = SentenceTransformerEmbeddingProvider(model_name="all-MiniLM-L6-v2")
        v2 = st_emb.embed_text("National Statistical Office Survey")
        assert len(v2) == 384
        assert st_emb.dimension == 384
    except Exception as e:
        pytest.skip(f"SentenceTransformer not available in current environment: {e}")


def test_txt_extractor_multi_encoding():
    """Verify TXT extractor handles UTF-8 and Latin-1 safely."""
    extractor = TXTTextExtractor()
    
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        f.write("MoSPI Consumer Price Index (CPI) 2026".encode("utf-8"))
        utf8_path = f.name

    try:
        res_utf8 = extractor.extract(utf8_path)
        assert "Consumer Price Index" in res_utf8[0]["text"]
        assert res_utf8[0]["page_number"] == 1
    finally:
        if os.path.exists(utf8_path):
            os.remove(utf8_path)

    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        f.write("Statistical Survey Référence Manual".encode("latin-1"))
        latin1_path = f.name

    try:
        res_latin1 = extractor.extract(latin1_path)
        assert "Statistical Survey" in res_latin1[0]["text"]
    finally:
        if os.path.exists(latin1_path):
            os.remove(latin1_path)


def test_pptx_extractor():
    """Verify PPTX text extractor handles slides and tables."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MoSPI Sampling Techniques"
    subtitle.text = "National Academy of Statistical Administration"

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        pptx_path = f.name

    try:
        extractor = PPTXTextExtractor()
        res = extractor.extract(pptx_path)
        assert len(res) >= 1
        assert "MoSPI Sampling Techniques" in res[0]["text"]
        assert res[0]["source_type"] == "slide"
        assert res[0]["page_number"] == 1
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)


def test_semantic_chunker():
    """Verify semantic chunking splits pages into overlapping chunks."""
    sample_text = (
        "1. Introduction to Official Sampling. "
        "Stratified random sampling is widely employed in socio-economic surveys conducted by NSO. "
        "The sampling frame is partitioned into mutually exclusive sub-populations called strata. "
        "2. Estimation Procedures. "
        "Unbiased estimators are derived by applying appropriate multiplier weights to sample values."
    )
    pages = [{"page_number": 1, "text": sample_text, "source_type": "page"}]
    chunks = DocumentProcessingService.chunk_pages(pages, target_word_count=15, overlap_word_count=5)
    
    assert len(chunks) >= 2
    for chunk in chunks:
        assert chunk["text"].strip()
        assert chunk["page_start"] == 1
        chunk_hash = DocumentProcessingService.compute_chunk_hash(uuid.uuid4(), 0, chunk["text"])
        assert len(chunk_hash) == 64


def test_sha256_file_hashing():
    """Verify SHA-256 file hashing is deterministic."""
    data = b"Official Statistics of India - 2026 Manual"
    hash1 = DocumentProcessingService.compute_sha256(data)
    hash2 = DocumentProcessingService.compute_sha256(data)
    assert hash1 == hash2
    assert len(hash1) == 64


def test_mcq_validator_quality_gate(db_session: Session):
    """Verify MCQValidator enforces options count, correct answer bounds, and source grounding."""
    fw = CompetencyFramework(name="AI_TEST_FW", description="Test FW")
    db_session.add(fw)
    db_session.commit()

    comp = Competency(framework_id=fw.id, name="Test Sampling", code="TEST_SAMP_99")
    db_session.add(comp)
    db_session.commit()

    source_text = "Stratified sampling ensures representation of key sub-groups in official surveys."

    # Valid MCQ
    valid_mcq = GeneratedMCQ(
        question="What does stratified sampling ensure in official surveys?",
        options=[
            GeneratedMCQOption(text="Representation of key sub-groups"),
            GeneratedMCQOption(text="Elimination of all sampling errors"),
            GeneratedMCQOption(text="Zero fieldwork costs"),
            GeneratedMCQOption(text="Arbitrary non-random selection")
        ],
        correct_answer=0,
        explanation="Stratified sampling ensures representation of key sub-groups as documented.",
        competency_code="TEST_SAMP_99",
        difficulty="MEDIUM",
        source_chunk_ids=[uuid.uuid4()]
    )
    is_valid, reasons, score = MCQValidator.validate(db_session, valid_mcq, source_text)
    assert is_valid is True
    assert score >= 0.2

    # Invalid MCQ (less than 4 options)
    invalid_mcq = GeneratedMCQ(
        question="Incomplete options question?",
        options=[
            GeneratedMCQOption(text="Option A"),
            GeneratedMCQOption(text="Option B")
        ],
        correct_answer=0,
        explanation="Short",
        competency_code="TEST_SAMP_99",
        difficulty="EASY"
    )
    is_valid_inv, reasons_inv, _ = MCQValidator.validate(db_session, invalid_mcq, source_text)
    assert is_valid_inv is False
    assert any("4" in r for r in reasons_inv)


def test_role_diagnostic_generator(db_session: Session):
    """Verify RoleDiagnosticGenerator creates questions mapped to DB competencies with difficulty distribution."""
    fw = CompetencyFramework(name="STAT_DIAG_FW", description="Diag FW")
    db_session.add(fw)
    db_session.commit()

    c1 = Competency(framework_id=fw.id, name="Survey Methodology", code="METHOD_01")
    c2 = Competency(framework_id=fw.id, name="Data Quality & Audit", code="AUDIT_02")
    db_session.add_all([c1, c2])
    db_session.commit()

    role = JobRole(name="Statistical Investigator", code="ROLE_STAT_INV")
    db_session.add(role)
    db_session.commit()

    rc1 = RoleCompetency(job_role_id=role.id, competency_id=c1.id, required_level=3, weight=1.0)
    rc2 = RoleCompetency(job_role_id=role.id, competency_id=c2.id, required_level=2, weight=1.0)
    db_session.add_all([rc1, rc2])
    db_session.commit()

    diag = RoleDiagnosticGenerator.generate_role_assessment(
        db=db_session,
        job_role_id=role.id,
        total_questions=10
    )

    assert diag["job_role"] == "Statistical Investigator"
    assert diag["total_questions"] == 10
    assert len(diag["questions"]) == 10
    assert "competency_breakdown" in diag
    assert len(diag["competency_breakdown"]) == 2

    # Check difficulty distribution
    difficulties = [q["difficulty"] for q in diag["questions"]]
    assert "EASY" in difficulties
    assert "MEDIUM" in difficulties
    assert "HARD" in difficulties


def test_ai_health_endpoint(client: TestClient):
    """Verify /api/v1/ai/health returns expected health check structure."""
    resp = client.get("/api/v1/ai/health")
    assert resp.status_code == 200
    data = resp.json()
    assert "status" in data
    assert "llm" in data
    assert "embeddings" in data
    assert data["canonical_dimension"] == 384
