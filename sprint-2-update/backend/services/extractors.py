import abc
import os
import re
from typing import Dict, Any, List

class DocumentExtractor(abc.ABC):
    @abc.abstractmethod
    def extract(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        """
        Extracts content from a document and returns a normalized dictionary.
        Format:
        {
            "metadata": {"original_filename": str},
            "sections": [
                {
                    "index": int,
                    "text": str,
                    "metadata": {"page": int, "slide": int, "type": str}
                }
            ]
        }
        """
        pass

    def normalize_text(self, text: str) -> str:
        """Removes excessive whitespace and standardizes line endings"""
        if not text:
            return ""
        # Normalize newlines
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        # Remove multiple newlines
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Remove multiple spaces but preserve newlines
        text = re.sub(r'[ \t]+', ' ', text)
        return text.strip()


class PDFExtractor(DocumentExtractor):
    def extract(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        from pypdf import PdfReader
        
        normalized = {
            "metadata": {"original_filename": original_filename, "type": "pdf"},
            "sections": []
        }
        
        try:
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                raw_text = page.extract_text()
                clean_text = self.normalize_text(raw_text)
                if clean_text:
                    normalized["sections"].append({
                        "index": i + 1,
                        "text": clean_text,
                        "metadata": {"page": i + 1}
                    })
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
            
        return normalized

class DOCXExtractor(DocumentExtractor):
    def extract(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        from docx import Document
        
        normalized = {
            "metadata": {"original_filename": original_filename, "type": "docx"},
            "sections": []
        }
        
        try:
            doc = Document(file_path)
            # Group paragraphs roughly into chunks or treat the whole docx sequentially
            # For simplicity, we just chunk by large paragraphs or keep it simple.
            
            # Simple approach: Join paragraphs, chunk if too long, or just one big section 
            # for now let's just make each paragraph a section if it has content, or 
            # bundle them by 10 paragraphs to avoid excessive small sections.
            
            current_chunk = []
            chunk_index = 1
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    current_chunk.append(text)
                
                # Bundle every 10 paragraphs or so
                if len(current_chunk) >= 10:
                    joined_text = "\n\n".join(current_chunk)
                    normalized["sections"].append({
                        "index": chunk_index,
                        "text": self.normalize_text(joined_text),
                        "metadata": {"section": chunk_index}
                    })
                    chunk_index += 1
                    current_chunk = []
                    
            if current_chunk:
                joined_text = "\n\n".join(current_chunk)
                normalized["sections"].append({
                    "index": chunk_index,
                    "text": self.normalize_text(joined_text),
                    "metadata": {"section": chunk_index}
                })
                
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
            
        return normalized

class PPTXExtractor(DocumentExtractor):
    def extract(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        from pptx import Presentation
        
        normalized = {
            "metadata": {"original_filename": original_filename, "type": "pptx"},
            "sections": []
        }
        
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                joined_text = "\n".join(slide_text)
                clean_text = self.normalize_text(joined_text)
                
                if clean_text:
                    normalized["sections"].append({
                        "index": i + 1,
                        "text": clean_text,
                        "metadata": {"slide": i + 1}
                    })
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
            
        return normalized

class TXTExtractor(DocumentExtractor):
    def extract(self, file_path: str, original_filename: str) -> Dict[str, Any]:
        normalized = {
            "metadata": {"original_filename": original_filename, "type": "txt"},
            "sections": []
        }
        
        try:
            # Handle various encodings
            content = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
                    
            if content is None:
                raise Exception("Failed to decode text file with common encodings")
                
            clean_text = self.normalize_text(content)
            
            # For txt, just treat it as one section for now, or split by double newlines.
            if clean_text:
                normalized["sections"].append({
                    "index": 1,
                    "text": clean_text,
                    "metadata": {"section": 1}
                })
        except Exception as e:
            raise Exception(f"TXT extraction failed: {str(e)}")
            
        return normalized

class ExtractorFactory:
    @staticmethod
    def get_extractor(file_extension: str) -> DocumentExtractor:
        ext = file_extension.lower().lstrip('.')
        if ext == 'pdf':
            return PDFExtractor()
        elif ext == 'docx':
            return DOCXExtractor()
        elif ext == 'pptx':
            return PPTXExtractor()
        elif ext == 'txt':
            return TXTExtractor()
        else:
            raise ValueError(f"Unsupported file extension: {ext}")
